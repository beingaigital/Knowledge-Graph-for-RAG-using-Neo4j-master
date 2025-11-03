# -*- coding: utf-8 -*-
"""
主程序：基于 RAG（LlamaIndex 0.11.x + Chroma + Neo4j）与 LiteLLM 的六类产出生成器
"""
import os, json, argparse, datetime as dt, pathlib
from typing import Dict, Any, List, Tuple

import yaml
from litellm import completion
import matplotlib.pyplot as plt
from PIL import Image, ImageDraw

from llama_index.core import VectorStoreIndex
from llama_index.embeddings.huggingface import HuggingFaceEmbedding
try:
    from llama_index.vector_stores.chroma import ChromaVectorStore
except Exception:
    ChromaVectorStore = None
import chromadb
from neo4j import GraphDatabase

from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches, Pt

from templates.prompts import (
    A_GRAPHIC_BRIEF, B_VIDEO_SCRIPT, C_CAMPAIGN_PLAN,
    D_SHORTVIDEO_SCRIPT, E_XHS_NOTE, F_CRISIS_PLAN
)

def get_storage_context_with_chroma(persist_dir: str = "./chroma_db", collection_name: str = "pr_agent"):
    """Return (storage_context, used_chroma: bool). Falls back to in-memory if ChromaVectorStore is unavailable."""
    try:
        if ChromaVectorStore is None:
            raise ImportError("ChromaVectorStore plugin not installed")
        import chromadb  # ensure chromadb present
        client = chromadb.PersistentClient(path=persist_dir)
        collection = client.get_or_create_collection(collection_name)
        vs = ChromaVectorStore(chroma_collection=collection)
        return StorageContext.from_defaults(vector_store=vs), True
    except Exception as e:
        print(f"[WARN] Chroma unavailable or failed ({e}); falling back to in-memory vector store.")
        return StorageContext.from_defaults(), False


def ensure_dir(p: str):
    pathlib.Path(p).mkdir(parents=True, exist_ok=True)

def tsdir(base: str) -> str:
    t = dt.datetime.now().strftime("%Y%m%d_%H%M%S")
    out = os.path.join(base, t); ensure_dir(out); return out

def read_config() -> Dict[str, Any]:
    cfg_path = "config.yaml" if os.path.exists("config.yaml") else "config.example.yaml"
    return yaml.safe_load(open(cfg_path, "r", encoding="utf-8"))

class GraphRAG:
    def __init__(self, persist_dir: str, neo4j_uri: str, neo4j_user: str, neo4j_pwd: str, top_k: int = 10):
        self.client = chromadb.PersistentClient(path=persist_dir)
        self.coll = self.client.get_or_create_collection("pr_kb_v3")
        self.vector_store = ChromaVectorStore(chroma_collection=self.coll)
        self.embed_model = HuggingFaceEmbedding(model_name="BAAI/bge-m3")
        self.index = VectorStoreIndex.from_vector_store(self.vector_store, embed_model=self.embed_model, show_progress=False)
        self.retriever = self.index.as_retriever(similarity_top_k=top_k)
        self.driver = GraphDatabase.driver(neo4j_uri, auth=(neo4j_user, neo4j_pwd))

    def close(self):
        self.driver.close()

    def retrieve(self, query: str, k: int = 10) -> List[Dict[str, Any]]:
        nodes = self.retriever.retrieve(query)[:k]
        out = []
        for i, n in enumerate(nodes, 1):
            meta = getattr(n.node, "metadata", {}) or {}
            out.append({"text": n.node.get_content(), "score": float(getattr(n, "score", 0.0)), "meta": meta})
        return out

    def fetch_graph(self, goal: str) -> Dict[str, List[str]]:
        cypher = """
        OPTIONAL MATCH (g:PRGoal {name:$goal})-[:INCLUDES]->(s:Strategy)
        OPTIONAL MATCH (s)-[:USES]->(c:Channel)
        OPTIONAL MATCH (s)-[:ILLUSTRATED_BY]->(cs:CaseStudy)
        OPTIONAL MATCH (s)-[:APPLIES_TO]->(p:Persona)
        RETURN collect(DISTINCT s.title) AS strategies,
               collect(DISTINCT c.name) AS channels,
               collect(DISTINCT cs.title) AS cases,
               collect(DISTINCT p.name) AS personas
        """
        with self.driver.session() as s:
            rec = s.run(cypher, goal=goal).single()
            if not rec:
                return {"strategies": [], "channels": [], "cases": [], "personas": []}
            return {k: [x for x in rec[k] if x] for k in ["strategies","channels","cases","personas"]}

def llm_complete(provider: str, model: str, prompt: str, max_tokens=2048, temperature=0.6) -> str:
    model_id = model if "/" in model else f"{provider}/{model}"
    rsp = completion(model=model_id, messages=[{"role":"user","content":prompt}], max_tokens=max_tokens, temperature=temperature)
    try:
        return rsp.choices[0].message["content"]
    except Exception:
        import json as _json
        return _json.dumps(rsp, ensure_ascii=False)

def save_graphics_placeholders(out_dir: str, campaign_name: str, count: int = 3):
    w, h = 3508, 4961
    for i in range(1, count+1):
        img = Image.new("RGB", (w, h), (240, 240, 240))
        d = ImageDraw.Draw(img)
        d.multiline_text((w//6, h//2-50), f"{campaign_name}\nA3/300dpi Placeholder\n#{i}", fill=(40,40,40), spacing=10)
        img.save(os.path.join(out_dir, f"A_graphic_{i:02d}.jpg"), "JPEG", quality=95)

def plot_budget_pie(budgets: Dict[str, float], save_path: str):
    labels = list(budgets.keys()); sizes = list(budgets.values())
    plt.figure(figsize=(6,6)); plt.pie(sizes, labels=labels, autopct='%1.1f%%', startangle=90); plt.axis('equal')
    plt.title("预算分配"); plt.savefig(save_path, dpi=200, bbox_inches='tight'); plt.close()

def plot_gantt(schedule: List[Tuple[str,int]], save_path: str):
    plt.figure(figsize=(8,4)); left = 0
    for name, weeks in schedule:
        plt.barh([0], [weeks], left=[left]); left += weeks; plt.text(left - weeks/2, 0, name, ha='center', va='center')
    plt.yticks([]); plt.xlabel("周"); plt.title("执行甘特图"); plt.savefig(save_path, dpi=200, bbox_inches='tight'); plt.close()

def export_word_plan(path: str, title: str, outline: str, budget_png: str, gantt_png: str):
    doc = Document(); doc.add_heading(title, 0); doc.add_paragraph(outline)
    doc.add_heading("预算分配", level=1); doc.add_picture(budget_png, width=Inches(5))
    doc.add_heading("执行甘特图", level=1); doc.add_picture(gantt_png, width=Inches(6)); doc.save(path)

def export_ppt_plan(path: str, title: str, outline_points: List[str], budget_png: str, gantt_png: str):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0]); slide.shapes.title.text = title; slide.placeholders[1].text = dt.datetime.now().strftime("%Y-%m-%d")
    slide = prs.slides.add_slide(prs.slide_layouts[1]); slide.shapes.title.text = "方案大纲"; tf = slide.placeholders[1].text_frame; tf.clear()
    for pt in outline_points: p = tf.add_paragraph(); p.text = pt; p.level = 0
    for sec_title, img in [("预算分配", budget_png), ("执行甘特图", gantt_png)]:
        slide = prs.slides.add_slide(prs.slide_layouts[5]); slide.shapes.title.text = sec_title; slide.shapes.add_picture(img, Inches(1), Inches(1.5), height=Inches(4.5))
    prs.save(path)

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--enterprise_name", required=True)
    ap.add_argument("--enterprise_stage", required=True, choices=["初创企业","中小微企业","大型国企央企","政府事业单位","NGO公益组织"])
    ap.add_argument("--industry", required=True)
    ap.add_argument("--market_type", required=True, choices=["ToB","ToC","ToG"])
    ap.add_argument("--pr_goal", required=True, choices=["品牌认知","产品推广","市场拓展","危机公关","用户心智","政府关系","社区及用户运营","形象提升"])
    ap.add_argument("--pr_cycle", required=True)
    ap.add_argument("--pr_budget", required=True)
    ap.add_argument("--innovation", required=True, choices=["保守","适度创新","不拘一格"])
    ap.add_argument("--outputs", default="A,B,C,D,E,F")
    args = ap.parse_args()

    cfg = read_config()
    out_dir = tsdir(cfg["paths"]["output_dir"])

    rag = GraphRAG(
        persist_dir=cfg["paths"]["persist_dir"],
        neo4j_uri=cfg["neo4j"]["uri"],
        neo4j_user=cfg["neo4j"]["user"],
        neo4j_pwd=cfg["neo4j"]["password"],
        top_k=cfg["retrieval"]["top_k"]
    )

    query = f"{args.enterprise_stage} {args.industry} {args.market_type} 目标:{args.pr_goal} 创新:{args.innovation}"
    vec_hits = rag.retrieve(query, k=cfg["retrieval"]["top_k"])
    g = rag.fetch_graph(args.pr_goal)

    ctx_parts = []
    for i, h in enumerate(vec_hits, 1):
        src = h["meta"].get("source", "") if isinstance(h["meta"], dict) else ""
        ctx_parts.append(f"[{i}] {h['text'][:800]}\n— 来源：{src}")
    graph_part = f"策略: {g.get('strategies', [])}\n渠道: {g.get('channels', [])}\n案例: {g.get('cases', [])}\n人群: {g.get('personas', [])}"
    context = "\n\n".join(ctx_parts + [graph_part])[: cfg["retrieval"]["max_context_chars"]]

    vars_text = json.dumps({
        "企业名称": args.enterprise_name,
        "企业类型": {"阶段": args.enterprise_stage, "行业": args.industry, "市场": args.market_type},
        "公关目标": args.pr_goal, "公关周期": args.pr_cycle, "公关预算": args.pr_budget, "创新程度": args.innovation
    }, ensure_ascii=False)

    provider = cfg["llm"]["provider"]; model = cfg["llm"]["model"]
    max_tokens = int(cfg["llm"].get("max_tokens", 2048)); temperature = float(cfg["llm"].get("temperature", 0.6))
    want = [x.strip().upper() for x in args.outputs.split(",") if x.strip()]

    # A
    if "A" in want:
        a_dir = os.path.join(out_dir, "A_graphics"); ensure_dir(a_dir)
        brief = llm_complete(provider, model, A_GRAPHIC_BRIEF.format(context=context, vars=vars_text), max_tokens, temperature)
        open(os.path.join(a_dir, "A_brief.md"), "w", encoding="utf-8").write(brief)
        save_graphics_placeholders(a_dir, args.enterprise_name, 3)

    # B
    if "B" in want:
        b_dir = os.path.join(out_dir, "B_corp_video"); ensure_dir(b_dir)
        script = llm_complete(provider, model, B_VIDEO_SCRIPT.format(context=context, vars=vars_text), max_tokens, temperature)
        open(os.path.join(b_dir, "B_script_shotlist.md"), "w", encoding="utf-8").write(script)

    # C
    if "C" in want:
        c_dir = os.path.join(out_dir, "C_campaign_plan"); ensure_dir(c_dir)
        outline = llm_complete(provider, model, C_CAMPAIGN_PLAN.format(context=context, vars=vars_text), max_tokens, temperature)
        budgets = {"品牌传播":40, "内容制作":35, "投放":20, "监测评估":5}
        budget_png = os.path.join(c_dir, "budget.png"); plot_budget_pie(budgets, budget_png)
        gantt = [("预热", 2), ("爆发", 6), ("延续", 6), ("复盘", 2)]
        gantt_png = os.path.join(c_dir, "gantt.png"); plot_gantt(gantt, gantt_png)
        export_word_plan(os.path.join(c_dir, "campaign_plan.docx"), f"{args.enterprise_name} 公关营销策划案", outline, budget_png, gantt_png)
        outline_points = [ln.strip("-• ").strip() for ln in outline.splitlines() if ln.strip() and len(ln) < 80][:12]
        export_ppt_plan(os.path.join(c_dir, "campaign_plan.pptx"), f"{args.enterprise_name} 公关营销策划案", outline_points, budget_png, gantt_png)
        open(os.path.join(c_dir, "campaign_plan_outline.md"), "w", encoding="utf-8").write(outline)

    # D
    if "D" in want:
        d_dir = os.path.join(out_dir, "D_shortvideo"); ensure_dir(d_dir)
        sc = llm_complete(provider, model, D_SHORTVIDEO_SCRIPT.format(context=context, vars=vars_text), max_tokens, temperature)
        open(os.path.join(d_dir, "D_shortvideo_scripts.md"), "w", encoding="utf-8").write(sc)

    # E
    if "E" in want:
        e_dir = os.path.join(out_dir, "E_xiaohongshu"); ensure_dir(e_dir)
        note = llm_complete(provider, model, E_XHS_NOTE.format(context=context, vars=vars_text), max_tokens, temperature)
        open(os.path.join(e_dir, "E_note.md"), "w", encoding="utf-8").write(note)

    # F
    if "F" in want:
        f_dir = os.path.join(out_dir, "F_crisis_plan"); ensure_dir(f_dir)
        outline = llm_complete(provider, model, F_CRISIS_PLAN.format(context=context, vars=vars_text), max_tokens, temperature)
        budgets = {"监测与分析":25, "媒体与社区沟通":35, "内容制作":25, "培训演练":15}
        budget_png = os.path.join(f_dir, "budget.png"); plot_budget_pie(budgets, budget_png)
        gantt = [("第一响应", 1), ("沟通与澄清", 2), ("修复与重建", 4), ("复盘优化", 1)]
        gantt_png = os.path.join(f_dir, "gantt.png"); plot_gantt(gantt, gantt_png)
        export_word_plan(os.path.join(f_dir, "crisis_plan.docx"), f"{args.enterprise_name} 危机公关应对方案", outline, budget_png, gantt_png)
        outline_points = [ln.strip("-• ").strip() for ln in outline.splitlines() if ln.strip() and len(ln) < 80][:12]
        export_ppt_plan(os.path.join(f_dir, "crisis_plan.pptx"), f"{args.enterprise_name} 危机公关应对方案", outline_points, budget_png, gantt_png)
        open(os.path.join(f_dir, "crisis_plan_outline.md"), "w", encoding="utf-8").write(outline)

    rag.close()
    print("✅ 生成完成，输出目录：", out_dir)

if __name__ == "__main__":
    main()
